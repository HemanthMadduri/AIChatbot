from flask import Flask, request, jsonify, Response, stream_with_context
from flask_cors import CORS
import requests
from dotenv import load_dotenv
import os
from werkzeug.utils import secure_filename
import base64
import json


try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import cv2
except ImportError:
    cv2 = None

try:
    from groq import Groq
except ImportError:
    Groq = None

load_dotenv()

# Configuration: Choose API provider (groq or ollama)
API_PROVIDER = os.getenv("API_PROVIDER", "groq").lower()  # Default to Groq for cloud deployment

# Groq Configuration (free cloud API)
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")
GROQ_MODEL = os.getenv("GROQ_MODEL", "llama-3.3-70b-versatile")  # Fast and capable

# Ollama Configuration (local installation)
OLLAMA_URL = os.getenv("OLLAMA_HOST", "http://127.0.0.1:11434")
OLLAMA_API_URL = OLLAMA_URL + "/api/generate"
OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "llava")

MAX_TOKENS = 512  # lower = faster response

app = Flask(__name__)
CORS(app)


UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'png', 'jpg', 'jpeg', 'gif', 'mp4', 'avi', 'mov', 'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx'}
IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}
VIDEO_EXTENSIONS = {'mp4', 'avi', 'mov'}
DOCUMENT_EXTENSIONS = {'txt', 'pdf', 'doc', 'docx', 'ppt', 'pptx'}

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
# 200 MB max for requests (videos can be large)
app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024

@app.errorhandler(413)
def request_entity_too_large(error):
    return jsonify({
        "error": "File too large",
        "message": "The uploaded file exceeds the size limit (200 MB). Try a shorter/smaller video or compress it."
    }), 413

messages = [
    {"role": "system", "content": "You are a helpful AI assistant. Respond in single line if possible. When analyzing images or video frames, describe what you see in detail. When the user attaches documents (PDF, Word, PowerPoint), you will receive the extracted text content—answer questions about that content. When you receive multiple images as video frames, summarize the video; do not say you cannot view or analyze video."}
]

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_file_extension(filename):
    return filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

@app.route("/health", methods=["GET"])
def health():
    try:
        if API_PROVIDER == "groq":
            if not GROQ_API_KEY:
                return jsonify({
                    "status": "error",
                    "message": "GROQ_API_KEY not configured"
                }), 500
            return jsonify({
                "status": "ok",
                "provider": "groq",
                "model": GROQ_MODEL
            })
        else:  # ollama
            response = requests.get(f"{OLLAMA_URL}/api/tags", timeout=5)
            if response.status_code == 200:
                models = response.json()
                return jsonify({
                    "status": "ok",
                    "provider": "ollama",
                    "ollama_url": OLLAMA_URL,
                    "models": models
                })
            else:
                return jsonify({"status": "error", "message": f"Ollama returned {response.status_code}"}), 500
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": f"Error: {str(e)}"
        }), 500

def encode_image_to_base64(filepath):
    try:
        with open(filepath, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    except Exception as e:
        print(f"Error encoding image: {e}")
        return None


def extract_frames_from_video(filepath, max_frames=5, interval_sec=3, max_width=480):
    """Extract sampled frames from video as base64 JPEGs for vision model (tuned for speed)."""
    if cv2 is None:
        return []
    try:
        cap = cv2.VideoCapture(filepath)
        if not cap.isOpened():
            return []
        fps = cap.get(cv2.CAP_PROP_FPS) or 1
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
        if total_frames <= 0:
            return []
        frame_interval = max(1, int(fps * interval_sec))
        frames_base64 = []
        frame_idx = 0
        jpeg_params = [cv2.IMWRITE_JPEG_QUALITY, 75]  # smaller payload
        while len(frames_base64) < max_frames and frame_idx < total_frames:
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
            ret, frame = cap.read()
            if not ret or frame is None:
                break
            h, w = frame.shape[:2]
            if w > max_width:
                scale = max_width / w
                frame = cv2.resize(frame, (max_width, int(h * scale)))
            _, buf = cv2.imencode(".jpg", frame, jpeg_params)
            frames_base64.append(base64.b64encode(buf.tobytes()).decode("utf-8"))
            frame_idx += frame_interval
        cap.release()
        return frames_base64
    except Exception as e:
        print(f"Error extracting video frames: {e}")
        return []


def extract_text_from_document(filepath):
    try:
        extension = get_file_extension(filepath)
        
        if extension == 'txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        
        elif extension == 'pdf':
            if PyPDF2 is None:
                return "[PDF file - PyPDF2 not installed. Run: pip install PyPDF2]"
            try:
                text = ""
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text()
                return text if text else "[PDF file - could not extract text]"
            except Exception as pdf_error:
                return f"[Error reading PDF: {str(pdf_error)}]"
        
        elif extension in ['doc', 'docx']:
            if Document is None:
                return "[DOCX file - python-docx not installed. Run: pip install python-docx]"
            try:
                doc = Document(filepath)
                return '\n'.join([para.text for para in doc.paragraphs])
            except Exception as docx_error:
                return f"[Error reading DOCX: {str(docx_error)}]"
        
        elif extension == 'pptx':
            if Presentation is None:
                return "[PPTX file - python-pptx not installed. Run: pip install python-pptx]"
            try:
                prs = Presentation(filepath)
                parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    slide_texts.append(t)
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                                if row_text:
                                    slide_texts.append(" | ".join(row_text))
                    if slide_texts:
                        parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
                return "\n\n".join(parts) if parts else "[PPTX file - no text content found]"
            except Exception as pptx_error:
                return f"[Error reading PPTX: {str(pptx_error)}]"
        
        elif extension == 'ppt':
            return "[Legacy .ppt format - text extraction not supported. Please save the presentation as .pptx and upload again.]"
        
        return f"[File: {filepath}]"
    except Exception as e:
        return f"[Error reading file: {str(e)}]"


def call_groq_api(user_message_content, has_images=False):
    """Call Groq API for chat completion"""
    if Groq is None:
        return "Error: Groq library not installed. Run: pip install groq"
    
    if not GROQ_API_KEY:
        return "Error: GROQ_API_KEY not configured. Get a free API key from https://console.groq.com"
    
    try:
        client = Groq(api_key=GROQ_API_KEY)
        
        # Note: Groq doesn't support image inputs with text models
        # For image analysis, you'd need a vision model or describe that images aren't supported
        if has_images:
            response_text = "Note: This deployment uses Groq API which doesn't support image analysis. "
            response_text += "For image/video support, you need to run Ollama locally with llava model. "
            response_text += f"\n\nRegarding your message: {user_message_content}"
        else:
            # Build conversation for Groq
            groq_messages = []
            for msg in messages:
                role = "assistant" if msg["role"] == "assistant" else "user"
                if msg["role"] == "system":
                    role = "system"
                groq_messages.append({"role": role, "content": msg["content"]})
            
            chat_completion = client.chat.completions.create(
                messages=groq_messages,
                model=GROQ_MODEL,
                temperature=0.7,
                max_tokens=MAX_TOKENS * 2,  # Groq is fast, we can afford more tokens
            )
            response_text = chat_completion.choices[0].message.content
        
        return response_text
    except Exception as e:
        return f"Error calling Groq API: {str(e)}"

@app.route("/chat", methods=["POST"])
def chat():
    try:
        if request.is_json:
            data = request.get_json()
            user_message = data.get("message", "")
            file_info = []
            image_base64 = None
            video_frames = []
            video_path = None
        else:
            user_message = request.form.get("message", "")
            file_info = []
            image_base64 = None
            video_frames = []
            video_path = None
            
            for key in request.files:
                file = request.files[key]
                if file and allowed_file(file.filename):
                    filename = secure_filename(file.filename)
                    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                    file.save(filepath)
                    
                    extension = get_file_extension(filename)
                    
                    if extension in IMAGE_EXTENSIONS:
                        image_base64 = encode_image_to_base64(filepath)
                        if image_base64:
                            file_info.append(f"[Image attached: {filename}]")
                    
                    elif extension in VIDEO_EXTENSIONS:
                        file_info.append(f"[Video attached: {filename}]")
                        video_path = filepath
                    
                    elif extension in DOCUMENT_EXTENSIONS:
                        doc_text = extract_text_from_document(filepath)
                        # Allow more content for presentations (e.g. multiple slides)
                        max_chars = 8000 if extension in ('ppt', 'pptx') else 5000
                        file_info.append(f"[Document attached: {filename}]\nContent:\n{doc_text[:max_chars]}")
                    
                    else:
                        file_info.append(f"File uploaded: {filename}")
            
            if video_path:
                video_frames = extract_frames_from_video(video_path, max_frames=5, interval_sec=3)
                if not video_frames and cv2 is None:
                    return jsonify({
                        "reply": "Video analysis requires opencv. Install it with: pip install opencv-python-headless"
                    })
                elif not video_frames:
                    return jsonify({
                        "reply": "Could not extract frames from the video. Check that the file is a valid video (e.g. MP4) and not corrupted."
                    })
        
        full_message = user_message
        if file_info:
            full_message += "\n" + "\n".join(file_info)
        
        if not full_message.strip():
            return jsonify({"error": "Message or file is required"}), 400

        messages.append({"role": "user", "content": full_message})
        
        # Use Groq API if configured, otherwise Ollama
        if API_PROVIDER == "groq":
            has_media = bool(image_base64 or video_frames)
            reply = call_groq_api(full_message, has_images=has_media)
            messages.append({"role": "assistant", "content": reply})
            return jsonify({"reply": reply, "files_processed": len(file_info), "provider": "groq"})
        
        # Ollama implementation (original code)
        prompt_text = ""
        for msg in messages:
            role = msg["role"].capitalize()
            content = msg["content"]
            prompt_text += f"{role}: {content}\n"
        prompt_text += "Assistant:"
        
        use_stream = request.form.get("stream") == "1"
        if video_frames:
            video_prompt = (
                "These images are frames from a video, in chronological order. "
                "Summarize what happens in the video: describe the main actions, people, and events you see. "
            )
            if user_message and user_message.strip():
                video_prompt += f"\n\nUser question: {user_message}"
            else:
                video_prompt += "Give a concise summary."
            payload = {
                "model": OLLAMA_MODEL,
                "prompt": video_prompt,
                "images": video_frames,
                "temperature": 0.7,
                "max_tokens": MAX_TOKENS,
                "stream": use_stream
            }
        elif image_base64:
            payload = {
                "model": OLLAMA_MODEL,
                "prompt": f"Please analyze this image and answer: {user_message if user_message else 'What is in this image?'}",
                "images": [image_base64],
                "temperature": 0.7,
                "max_tokens": MAX_TOKENS,
                "stream": use_stream
            }
        else:
            payload = {
                "model": OLLAMA_MODEL,
                "prompt": prompt_text,
                "temperature": 0.7,
                "max_tokens": MAX_TOKENS,
                "stream": use_stream
            }
        
        req_timeout = 90 if video_frames else 60
        print(f"Sending request to {OLLAMA_API_URL} with model {OLLAMA_MODEL} stream={use_stream}")
        
        if use_stream:
            resp = requests.post(OLLAMA_API_URL, json=payload, timeout=req_timeout, stream=True)
            if resp.status_code != 200:
                reply = f"Error: Ollama returned status {resp.status_code}. Make sure Ollama is running and the '{OLLAMA_MODEL}' model is installed."
                messages.append({"role": "assistant", "content": reply})
                return jsonify({"reply": reply, "files_processed": len(file_info)})
            def generate():
                full_reply = []
                for line in resp.iter_lines(decode_unicode=True):
                    if not line:
                        continue
                    try:
                        chunk = json.loads(line)
                        delta = chunk.get("response", "")
                        if delta:
                            full_reply.append(delta)
                            yield f"data: {json.dumps({'delta': delta})}\n\n"
                        if chunk.get("done"):
                            break
                    except (json.JSONDecodeError, KeyError):
                        continue
                reply = "".join(full_reply) if full_reply else "No response from Ollama"
                messages.append({"role": "assistant", "content": reply})
                yield f"data: {json.dumps({'done': True, 'reply': reply})}\n\n"
            return Response(
                stream_with_context(generate()),
                content_type="text/event-stream",
                headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"}
            )
        
        response = requests.post(OLLAMA_API_URL, json=payload, timeout=req_timeout)
        print(f"Response status: {response.status_code}")
        reply = response.json().get("response", "No response from Ollama")
        if response.status_code != 200:
            print(f"Ollama API Error: {response.status_code} - {response.text}")
            reply = f"Error: Ollama returned status {response.status_code}. Make sure Ollama is running and the '{OLLAMA_MODEL}' model is installed."
        messages.append({"role": "assistant", "content": reply})
        return jsonify({"reply": reply, "files_processed": len(file_info), "provider": "ollama"})

    except Exception as e:
        print(f"Error in /chat: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    app.run(debug=True, threaded=True)